from fastapi import APIRouter, Depends, HTTPException, Form
from sqlalchemy import text
from sqlalchemy.orm import Session
from utils import get_db, check_permission
from verify_token import get_current_user
import os
import google.generativeai as genai
import base64

# Optional extractors
try:
    from pypdf import PdfReader  # pdf
except Exception:
    PdfReader = None
try:
    import docx  # python-docx
except Exception:
    docx = None
try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None
try:
    import openpyxl  # xlsx
except Exception:
    openpyxl = None

router = APIRouter()

genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
_model_name = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")

TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".log", ".js", ".ts", ".jsx", ".tsx", ".py", ".java", ".c", ".cpp", ".go", ".rb", ".php", ".sh", ".html", ".css"}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}
PPTX_EXTS = {".ppt", ".pptx"}
XLSX_EXTS = {".xls", ".xlsx"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}
AUDIO_EXTS = {".mp3", ".wav", ".ogg", ".m4a"}
VIDEO_EXTS = {".mp4", ".webm", ".ogg", ".mov"}
MAX_BYTES = 200_000


def _read_text_file(path: str) -> str:
    try:
        with open(path, "rb") as f:
            data = f.read(MAX_BYTES)
        try:
            return data.decode("utf-8")
        except Exception:
            return data.decode("latin-1", errors="ignore")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to read file: {str(e)}")


def _extract_pdf_text(path: str) -> str:
    if not PdfReader:
        raise HTTPException(status_code=500, detail="PDF support requires 'pypdf'. Add it to requirements and install.")
    try:
        reader = PdfReader(path)
        parts = []
        total = 0
        for page in reader.pages:
            txt = page.extract_text() or ""
            if not txt:
                continue
            b = txt.encode("utf-8")
            if total + len(b) > MAX_BYTES:
                b = b[: MAX_BYTES - total]
            parts.append(b.decode("utf-8", errors="ignore"))
            total += len(b)
            if total >= MAX_BYTES:
                break
        return "\n\n".join(parts)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF parse failed: {str(e)}")


def _extract_docx_text(path: str) -> str:
    if not docx:
        raise HTTPException(status_code=500, detail="DOCX support requires 'python-docx'. Add it to requirements and install.")
    try:
        d = docx.Document(path)
        parts = []
        total = 0
        for p in d.paragraphs:
            t = p.text or ""
            if not t:
                continue
            b = (t + "\n").encode("utf-8")
            if total + len(b) > MAX_BYTES:
                b = b[: MAX_BYTES - total]
            parts.append(b.decode("utf-8", errors="ignore"))
            total += len(b)
            if total >= MAX_BYTES:
                break
        return "".join(parts)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"DOCX parse failed: {str(e)}")


def _extract_pptx_text(path: str) -> str:
    if not Presentation:
        raise HTTPException(status_code=500, detail="PPT/PPTX support requires 'python-pptx'. Add it to requirements and install.")
    try:
        prs = Presentation(path)
        parts = []
        total = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in p.runs)
                        if not t:
                            continue
                        b = (t + "\n").encode("utf-8")
                        if total + len(b) > MAX_BYTES:
                            b = b[: MAX_BYTES - total]
                        parts.append(b.decode("utf-8", errors="ignore"))
                        total += len(b)
                        if total >= MAX_BYTES:
                            break
            if total >= MAX_BYTES:
                break
        return "".join(parts)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX parse failed: {str(e)}")


def _extract_xlsx_text(path: str) -> str:
    if not openpyxl:
        raise HTTPException(status_code=500, detail="XLS/XLSX support requires 'openpyxl'. Add it to requirements and install.")
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        total = 0
        for ws in wb.worksheets[:5]:  # cap sheets
            for row in ws.iter_rows(min_row=1, max_row=200, values_only=True):  # cap rows
                vals = [str(v) for v in row if v is not None]
                if not vals:
                    continue
                t = ", ".join(vals) + "\n"
                b = t.encode("utf-8")
                if total + len(b) > MAX_BYTES:
                    b = b[: MAX_BYTES - total]
                parts.append(b.decode("utf-8", errors="ignore"))
                total += len(b)
                if total >= MAX_BYTES:
                    break
            if total >= MAX_BYTES:
                break
        return "".join(parts)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"XLSX parse failed: {str(e)}")


def _read_bytes(path: str) -> bytes:
    try:
        with open(path, "rb") as f:
            return f.read(min(MAX_BYTES, os.path.getsize(path)))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to read file bytes: {str(e)}")


@router.post("/summarize")
def summarize(file_id: int = Form(...), db: Session = Depends(get_db), current_user: dict = Depends(get_current_user)):
    user_id = current_user["user_id"]
    perm = check_permission(db, user_id, file_id=file_id, operation="view")
    if not perm:
        raise HTTPException(status_code=403, detail="No permission")
    row = db.execute(text("SELECT file_path, file_name FROM files WHERE file_id = :fid"), {"fid": file_id}).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="File not found")
    file_path = row.file_path
    _, ext = os.path.splitext(file_path or "")
    ext = ext.lower()

    # Try to extract or attach content depending on type
    content = None
    image_part = None
    if ext in TEXT_EXTS:
        content = _read_text_file(file_path)
    elif ext in PDF_EXTS:
        content = _extract_pdf_text(file_path)
    elif ext in DOCX_EXTS:
        content = _extract_docx_text(file_path)
    elif ext in PPTX_EXTS:
        content = _extract_pptx_text(file_path)
    elif ext in XLSX_EXTS:
        content = _extract_xlsx_text(file_path)
    elif ext in IMAGE_EXTS:
        raw = _read_bytes(file_path)
        mime = f"image/{'jpeg' if ext in {'.jpg', '.jpeg'} else ext.lstrip('.')}"
        image_part = {"mime_type": mime, "data": base64.b64encode(raw).decode("ascii")}
    else:
        # Fallback: send bytes as attachment with a generic description
        raw = _read_bytes(file_path)
        image_part = {"mime_type": "application/octet-stream", "data": base64.b64encode(raw).decode("ascii")}

    if not os.getenv("GEMINI_API_KEY"):
        raise HTTPException(status_code=500, detail="GEMINI_API_KEY not set")
    model = genai.GenerativeModel(_model_name)
    parts = [
        "You are a concise assistant. Summarize the following file for a general user. Return: title, short summary (5-8 lines), key points as bullets.",
        f"Filename: {row.file_name}",
    ]
    if content:
        parts.append("\nContent:\n" + content)
    if image_part:
        parts.append(image_part)
    try:
        resp = model.generate_content(parts)
        text_out = resp.text if hasattr(resp, "text") else str(resp)
        return {"file_id": file_id, "summary": text_out}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gemini error: {str(e)}")


@router.post("/ask")
def ask(question: str = Form(...), file_id: int = Form(...), db: Session = Depends(get_db), current_user: dict = Depends(get_current_user)):
    user_id = current_user["user_id"]
    perm = check_permission(db, user_id, file_id=file_id, operation="view")
    if not perm:
        raise HTTPException(status_code=403, detail="No permission")
    row = db.execute(text("SELECT file_path, file_name FROM files WHERE file_id = :fid"), {"fid": file_id}).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="File not found")
    file_path = row.file_path
    _, ext = os.path.splitext(file_path or "")
    ext = ext.lower()

    content = None
    image_part = None
    if ext in TEXT_EXTS:
        content = _read_text_file(file_path)
    elif ext in PDF_EXTS:
        content = _extract_pdf_text(file_path)
    elif ext in DOCX_EXTS:
        content = _extract_docx_text(file_path)
    elif ext in PPTX_EXTS:
        content = _extract_pptx_text(file_path)
    elif ext in XLSX_EXTS:
        content = _extract_xlsx_text(file_path)
    elif ext in IMAGE_EXTS:
        raw = _read_bytes(file_path)
        mime = f"image/{'jpeg' if ext in {'.jpg', '.jpeg'} else ext.lstrip('.')}"
        image_part = {"mime_type": mime, "data": base64.b64encode(raw).decode("ascii")}
    else:
        raw = _read_bytes(file_path)
        image_part = {"mime_type": "application/octet-stream", "data": base64.b64encode(raw).decode("ascii")}

    if not os.getenv("GEMINI_API_KEY"):
        raise HTTPException(status_code=500, detail="GEMINI_API_KEY not set")
    model = genai.GenerativeModel(_model_name)
    parts = [
        "You are a helpful assistant. Answer the user question strictly using only the provided file. If not present, say you don't know. Provide short citation quotes.",
        f"Filename: {row.file_name}",
    ]
    if content:
        parts.append("\nFile Content:\n" + content)
    if image_part:
        parts.append(image_part)
    parts.append("\nQuestion: " + question)
    try:
        resp = model.generate_content(parts)
        text_out = resp.text if hasattr(resp, "text") else str(resp)
        return {"file_id": file_id, "answer": text_out}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gemini error: {str(e)}")
